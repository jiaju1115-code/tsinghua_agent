from __future__ import annotations

import shutil
from pathlib import Path
from typing import Any

from .common import dedupe_path, ensure_output_path
from .models import FilePlan


def _format_paragraph(paragraph: Any, *, size: int = 20, bold: bool = False, color: str = "24324A") -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    for run in paragraph.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)


def _add_notes(slide: Any, sources: list[dict[str, str]]) -> None:
    if not sources:
        return
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = "[Sources]\n" + "\n".join(
        f"- {source.get('title', '来源')}: {source.get('url', '')}" for source in sources
    )


def create_pptx(plan: FilePlan, output_path: str | Path | None = None, template_path: str | Path | None = None) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    path = dedupe_path(ensure_output_path(plan.title, "pptx", output_path))
    presentation = Presentation(str(template_path)) if template_path else Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    if not template_path:
        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = RGBColor.from_string("F6F4F8")
        title_slide.shapes.title.text = plan.title
        title_slide.placeholders[1].text = plan.subtitle or "清问·TsingAsk 可信校园事务智能体"
        title_para = title_slide.shapes.title.text_frame.paragraphs[0]
        _format_paragraph(title_para, size=50, bold=True, color="532A6E")
        title_para.alignment = PP_ALIGN.LEFT
        _format_paragraph(title_slide.placeholders[1].text_frame.paragraphs[0], size=24, color="596273")
        _add_notes(title_slide, plan.sources)
    slide_specs = plan.slides or [
        {"title": section.heading, "bullets": section.bullets or section.paragraphs, "table": section.table}
        for section in plan.sections
    ]
    for spec in slide_specs:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
        slide.shapes.title.text = str(spec.get("title", ""))
        _format_paragraph(slide.shapes.title.text_frame.paragraphs[0], size=35, bold=True, color="532A6E")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = list(spec.get("bullets", []))[:6]
        if not bullets and spec.get("table"):
            bullets = [" · ".join(str(value) for value in row) for row in spec["table"][:6]]
        for index, text in enumerate(bullets or ["[请根据实际情况补充]"]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(text)
            paragraph.level = 0
            paragraph.space_after = Pt(10)
            _format_paragraph(paragraph, size=20, color="24324A")
        _add_notes(slide, plan.sources)
    presentation.save(path)
    return path


def _replace_text_frame(frame: Any, replacements: dict[str, str]) -> int:
    changed = 0
    for paragraph in frame.paragraphs:
        original = paragraph.text
        updated = original
        for old, new in replacements.items():
            updated = updated.replace(old, new)
        if updated == original:
            continue
        if paragraph.runs:
            paragraph.runs[0].text = updated
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = updated
        changed += 1
    return changed


def modify_pptx(
    input_path: str | Path,
    *,
    replacements: dict[str, str] | None = None,
    output_path: str | Path | None = None,
) -> tuple[Path, int]:
    from pptx import Presentation

    source = Path(input_path).resolve()
    path = dedupe_path(ensure_output_path(f"{source.stem}_modified", "pptx", output_path))
    shutil.copy2(source, path)
    presentation = Presentation(path)
    changed = 0
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                changed += _replace_text_frame(shape.text_frame, replacements or {})
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        changed += _replace_text_frame(cell.text_frame, replacements or {})
    presentation.save(path)
    return path, changed


def read_pptx(path: str | Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(path)
    slides = []
    for number, slide in enumerate(presentation.slides, 1):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts.append(shape.text)
            if getattr(shape, "has_table", False):
                texts.extend(" | ".join(cell.text for cell in row.cells) for row in shape.table.rows)
        slides.append({"slide_number": number, "texts": texts})
    return {"format": "pptx", "slides": slides, "slide_count": len(slides)}
